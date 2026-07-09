#!/usr/bin/env python3
"""Convert a Markdown slide document into a simple 16:9 PPTX deck.

Input convention:
- `# Deck title` becomes the title slide heading.
- Each `## Slide N: Title` or `## Title` starts a new slide.
- The first non-empty paragraph after the slide heading becomes the subtitle when it is
  short and not a list/table/image.
- Markdown tables become PowerPoint tables.
- `- bullets`, `**bold lead:** text`, blockquotes, and Obsidian image embeds
  `![[image.png]]` are supported.
"""
from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MED_GRAY = RGBColor(0x94, 0xA3, 0xB8)

SLIDE_RE = re.compile(r"^##\s+(?:Slide\s+\d+\s*:\s*)?(.*)\s*$", re.I)
H1_RE = re.compile(r"^#\s+(.+)\s*$")
OBSIDIAN_IMAGE_RE = re.compile(r"!\[\[([^\]]+)\]\]")
MD_IMAGE_RE = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")


@dataclass
class SlideSpec:
    title: str
    lines: list[str] = field(default_factory=list)


def parse_markdown(path: Path) -> tuple[str, list[SlideSpec]]:
    deck_title = path.stem.replace("_", " ").title()
    slides: list[SlideSpec] = []
    current: SlideSpec | None = None

    for raw in path.read_text(encoding="utf-8").splitlines():
        line = raw.rstrip()
        h1 = H1_RE.match(line)
        slide = SLIDE_RE.match(line)
        if h1 and current is None and not slides:
            deck_title = h1.group(1).strip()
            continue
        if slide:
            current = SlideSpec(slide.group(1).strip() or "Untitled")
            slides.append(current)
            continue
        if current is not None:
            current.lines.append(line)

    if not slides:
        raise ValueError(f"No slides found in {path}. Use headings like '## Slide 1: Thesis'.")
    return deck_title, slides


def add_slide(prs: Presentation, title_text: str, subtitle_text: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    if subtitle_text:
        box2 = slide.shapes.add_textbox(Inches(0.7), Inches(0.82), Inches(11.9), Inches(0.75))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = clean_inline(subtitle_text)
        p2.font.size = Pt(26)
        p2.font.color.rgb = DARK
        p2.font.bold = True
    return slide


def clean_inline(text: str) -> str:
    text = text.strip()
    text = text.replace("**", "")
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text


def split_blocks(lines: list[str]) -> list[list[str]]:
    blocks: list[list[str]] = []
    cur: list[str] = []
    for line in lines:
        if not line.strip():
            if cur:
                blocks.append(cur)
                cur = []
        else:
            cur.append(line)
    if cur:
        blocks.append(cur)
    return blocks


def is_table(block: list[str]) -> bool:
    return len(block) >= 2 and "|" in block[0] and re.match(r"^\s*\|?\s*:?-{3,}:?", block[1]) is not None


def parse_table(block: list[str]) -> tuple[list[str], list[list[str]]]:
    def cells(row: str) -> list[str]:
        row = row.strip().strip("|")
        return [clean_inline(c.strip()) for c in row.split("|")]

    headers = cells(block[0])
    rows = [cells(r) for r in block[2:] if "|" in r]
    rows = [r + [""] * (len(headers) - len(r)) for r in rows]
    return headers, [r[: len(headers)] for r in rows]


def find_images(block: list[str], input_dir: Path, images_dir: Path | None) -> list[Path]:
    paths: list[Path] = []
    for line in block:
        refs = OBSIDIAN_IMAGE_RE.findall(line) + MD_IMAGE_RE.findall(line)
        for ref in refs:
            ref = ref.split("|")[0].strip()
            candidates = [input_dir / ref]
            if images_dir:
                candidates.append(images_dir / ref)
            candidates.append(Path(ref).expanduser())
            for c in candidates:
                if c.exists():
                    paths.append(c)
                    break
            else:
                print(f"Warning: image not found: {ref}")
    return paths


def add_table(slide, headers: list[str], rows: list[list[str]], top: float, left: float = 0.7, width: float = 11.9):
    if not rows:
        rows = [[""] * len(headers)]
    n_rows = len(rows) + 1
    n_cols = max(1, len(headers))
    row_h = min(0.55, max(0.28, 4.9 / max(n_rows, 1)))
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(row_h * n_rows))
    tbl = tbl_shape.table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
    return row_h * n_rows + 0.2


def add_text_block(slide, block: list[str], top: float, left: float = 0.7, width: float = 11.9, size: int = 15):
    height = min(5.6 - top, max(0.45, 0.32 * len(block) + 0.2))
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, line in enumerate(block):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        stripped = line.strip()
        is_quote = stripped.startswith("> ")
        is_bullet = stripped.startswith("- ") or stripped.startswith("* ")
        is_num = re.match(r"^\d+\.\s+", stripped) is not None
        if is_quote:
            stripped = "“" + stripped[2:].strip() + "”"
            p.font.italic = True
        elif is_bullet:
            stripped = stripped[2:].strip()
            p.level = 0
            p.text = "• " + clean_inline(stripped)
        elif is_num:
            p.text = clean_inline(stripped)
        else:
            p.text = clean_inline(stripped)
        if not p.text:
            p.text = clean_inline(stripped)
        p.font.size = Pt(size)
        p.font.color.rgb = GRAY if is_quote else DARK
        p.font.bold = stripped.startswith("**") or bool(re.match(r"^\*\*[^*]+:\*\*", line.strip()))
        p.space_after = Pt(4)
    return height + 0.12


def add_images(slide, images: Iterable[Path], top: float, left: float = 8.5, max_width: float = 4.0, max_height: float = 4.8):
    used = 0.0
    for img in images:
        try:
            slide.shapes.add_picture(str(img), Inches(left), Inches(top + used), width=Inches(max_width))
            used += min(max_height, 2.4) + 0.2
        except Exception as exc:  # noqa: BLE001
            print(f"Warning: could not add image {img}: {exc}")
    return used


def block_is_subtitle(block: list[str]) -> bool:
    if len(block) != 1:
        return False
    line = block[0].strip()
    if not line or line.startswith(("- ", "* ", "> ", "|", "!")):
        return False
    return len(clean_inline(line)) <= 115


def render_deck(input_md: Path, output_pptx: Path, images_dir: Path | None = None) -> None:
    _, specs = parse_markdown(input_md)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, spec in enumerate(specs, start=1):
        blocks = split_blocks(spec.lines)
        subtitle = None
        if blocks and block_is_subtitle(blocks[0]):
            subtitle = blocks.pop(0)[0]
        slide = add_slide(prs, f"Slide {idx}: {spec.title}", subtitle)
        top = 1.7 if subtitle else 1.25

        for block in blocks:
            images = find_images(block, input_md.parent, images_dir)
            text_lines = [l for l in block if not (OBSIDIAN_IMAGE_RE.search(l) or MD_IMAGE_RE.search(l))]
            if is_table(block):
                headers, rows = parse_table(block)
                top += add_table(slide, headers, rows, top)
            elif images and not text_lines:
                top += add_images(slide, images, top, left=0.8, max_width=11.0)
            else:
                if text_lines:
                    text_width = 7.4 if images else 11.9
                    top += add_text_block(slide, text_lines, top, width=text_width)
                if images:
                    add_images(slide, images, top=max(1.7, top - 1.2), left=8.7, max_width=3.8)
            if top > 6.35:
                print(f"Warning: slide may overflow: {spec.title}")
                break

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    print(f"Saved to {output_pptx}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Convert Markdown slides to PPTX.")
    parser.add_argument("input_md", type=Path, help="Markdown slide document, for example deep_dive_presentation.md")
    parser.add_argument("output_pptx", type=Path, nargs="?", help="Output .pptx path. Defaults to input stem + .pptx")
    parser.add_argument("--images-dir", type=Path, help="Directory for Obsidian image embeds")
    args = parser.parse_args()

    output = args.output_pptx or args.input_md.with_suffix(".pptx")
    render_deck(args.input_md.expanduser().resolve(), output.expanduser().resolve(), args.images_dir.expanduser().resolve() if args.images_dir else None)


if __name__ == "__main__":
    main()
